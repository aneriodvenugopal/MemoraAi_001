"""
RealApex Demo Video API Routes - SaaS Demo Content Generation
"""

from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Form
from fastapi.responses import Response, FileResponse
from typing import Optional, Literal, List
from pydantic import BaseModel
import os
import uuid
import json
import base64

router = APIRouter(prefix="/realapex-demos", tags=["RealApex Demos"])

# Import auth middleware
from middleware.auth import get_current_user


# Pydantic Models
class GenerateDemoScriptRequest(BaseModel):
    concept_title: str
    concept_subtitle: Optional[str] = ""
    category_name: Optional[str] = ""
    video_type: str
    video_type_description: Optional[str] = ""
    target_audience: str
    language: str
    custom_notes: Optional[str] = ""


class GeneratePresentationRequest(BaseModel):
    concept_title: str
    script: str
    language: str = "english"
    theme: str = "professional"  # professional, modern, minimal
    include_screenshots: bool = True
    screenshot_urls: Optional[List[str]] = []


class GenerateVoiceoverRequest(BaseModel):
    script: str
    voice: str = "nova"  # Now supports edge-tts voices
    speed: float = 1.0
    model: str = "edge-tts"  # Using free Edge TTS
    concept_title: Optional[str] = "voiceover"


def parse_script_to_slides(script: str, concept_title: str, num_screenshots: int = 0) -> list:
    """Parse script into slide structure WITHOUT using any AI - 100% FREE"""
    import re
    
    slides = []
    
    # Slide 1: Title
    slides.append({
        "slide_number": 1,
        "type": "title",
        "title": concept_title.replace(":", " -") if concept_title else "RealApex Demo",
        "subtitle": "Professional Real Estate SaaS Solution",
        "notes": "Welcome the audience"
    })
    
    # Parse script sections
    lines = script.strip().split('\n')
    current_bullets = []
    current_title = "Overview"
    slide_num = 2
    screenshot_count = 0
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Check for [SCREEN:...] markers - add screenshot slide
        if '[SCREEN:' in line or '[Screenshot:' in line.lower():
            # Save current bullets if any
            if current_bullets:
                slides.append({
                    "slide_number": slide_num,
                    "type": "content",
                    "title": current_title,
                    "bullets": current_bullets[:5],
                    "notes": ""
                })
                slide_num += 1
                current_bullets = []
            
            # Add screenshot slide
            caption = re.search(r'\[SCREEN:([^\]]+)\]', line, re.IGNORECASE)
            caption_text = caption.group(1) if caption else "Application Screenshot"
            slides.append({
                "slide_number": slide_num,
                "type": "screenshot",
                "title": f"Screenshot {screenshot_count + 1}",
                "caption": caption_text,
                "notes": "",
                "screenshot_index": screenshot_count
            })
            slide_num += 1
            screenshot_count += 1
            continue
        
        # Check for bullet points
        if line.startswith('-') or line.startswith('•') or line.startswith('*'):
            bullet = line.lstrip('-•* ').strip()
            if bullet:
                current_bullets.append(bullet)
        
        # Check for headers/titles
        elif line.startswith('#') or (len(line) < 50 and line.endswith(':')):
            # Save previous section
            if current_bullets:
                slides.append({
                    "slide_number": slide_num,
                    "type": "content",
                    "title": current_title,
                    "bullets": current_bullets[:5],
                    "notes": ""
                })
                slide_num += 1
                current_bullets = []
            current_title = line.lstrip('#').strip().rstrip(':')
        
        # Regular text - convert to bullet
        elif len(line) > 20:
            # Split long text into bullet points
            if len(line) < 100:
                current_bullets.append(line)
    
    # Add remaining bullets
    if current_bullets:
        slides.append({
            "slide_number": slide_num,
            "type": "content",
            "title": current_title,
            "bullets": current_bullets[:5],
            "notes": ""
        })
        slide_num += 1
    
    # Add screenshot slides for uploaded images if not already added
    while screenshot_count < num_screenshots:
        slides.append({
            "slide_number": slide_num,
            "type": "screenshot",
            "title": f"Feature Demo {screenshot_count + 1}",
            "caption": "Application screenshot",
            "notes": "",
            "screenshot_index": screenshot_count
        })
        slide_num += 1
        screenshot_count += 1
    
    # Add CTA slide
    slides.append({
        "slide_number": slide_num,
        "type": "cta",
        "title": "Get Started Today!",
        "bullets": ["Visit RealApex.in", "Request a Free Demo", "Contact Us"],
        "notes": "Call to action"
    })
    
    return slides


# Admin check helper
def check_admin(current_user: dict):
    """Check if user is admin or super_admin"""
    role = current_user.get("role", "")
    if role not in ["admin", "super_admin", "tenant_admin"]:
        raise HTTPException(status_code=403, detail="Admin access required")
    return True


@router.post("/generate-script")
async def generate_demo_script(
    request: GenerateDemoScriptRequest,
    current_user: dict = Depends(get_current_user)
):
    """Generate RealApex SaaS demo video script using Claude AI"""
    check_admin(current_user)
    
    try:
        from emergentintegrations.llm.chat import LlmChat, UserMessage
        import uuid
        
        EMERGENT_LLM_KEY = os.getenv('EMERGENT_LLM_KEY', '')
        
        # Build comprehensive prompt for SaaS demo
        prompt = f"""You are creating a professional YouTube demo video script for RealApex - a PropTech SaaS platform for Indian Real Estate developers and agents.

**VIDEO DETAILS:**
- Feature/Concept: {request.concept_title}
- Feature Description: {request.concept_subtitle}
- Category: {request.category_name}
- Video Type: {request.video_type} ({request.video_type_description})
- Target Audience: {request.target_audience}
- Language: {request.language}
{f"- Additional Notes: {request.custom_notes}" if request.custom_notes else ""}

**SCRIPT STRUCTURE (3-5 minutes video):**

1. **HOOK (15-20 seconds)**
   - Start with a pain point or challenge that {request.target_audience} face
   - Create curiosity about the solution

2. **INTRODUCTION (30 seconds)**
   - Introduce RealApex and this specific feature
   - Explain why this feature matters in 2026's competitive market

3. **FEATURE DEMONSTRATION (2-3 minutes)**
   - Step-by-step walkthrough of the feature
   - Highlight key UI elements and workflows
   - Use phrases like "As you can see on screen..." or "Notice how..."
   - Include specific benefits for {request.target_audience}

4. **KEY BENEFITS (30 seconds)**
   - List 3-4 concrete benefits
   - Include time savings, cost reduction, or efficiency gains
   - Mention ROI or competitive advantage

5. **CALL TO ACTION (20 seconds)**
   - Invite to try RealApex
   - Mention free demo or trial
   - Include website: RealApex.in

**STYLE GUIDELINES:**
- Use {request.language} naturally (if bilingual, mix smoothly)
- Professional but conversational tone
- Include pauses for screen transitions (mark as [PAUSE])
- Use phrases suitable for avatar/TTS narration
- Avoid complex technical jargon - explain simply
- Include enthusiasm appropriate for a product demo

**OUTPUT FORMAT:**
Write the complete script with clear section markers. Include [SCREEN: description] notes for what should appear on screen during that part.

Generate the script now:"""

        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=f"realapex-demo-{uuid.uuid4()}",
            system_message="You are an expert SaaS product marketing scriptwriter specializing in PropTech and Real Estate software demos for the Indian market."
        ).with_model("anthropic", "claude-sonnet-4-20250514")
        
        user_message = UserMessage(text=prompt)
        script = await chat.send_message(user_message)
        
        return {
            "success": True,
            "script": script,
            "concept_title": request.concept_title,
            "video_type": request.video_type,
            "target_audience": request.target_audience,
            "language": request.language
        }
        
    except Exception as e:
        print(f"❌ Demo script generation error: {str(e)}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/videos")
async def get_demo_videos(
    limit: int = 50,
    skip: int = 0,
    current_user: dict = Depends(get_current_user)
):
    """Get all generated RealApex demo videos"""
    check_admin(current_user)
    
    from motor.motor_asyncio import AsyncIOMotorClient
    
    MONGO_URL = os.getenv('MONGO_URL')
    DB_NAME = os.getenv('DB_NAME', 'test_database')
    client = AsyncIOMotorClient(MONGO_URL)
    db = client[DB_NAME]
    
    # Filter for RealApex Demo videos
    cursor = db.tutorai_generated_videos.find(
        {"subject": "RealApex Demo"}, 
        {"_id": 0}
    ).sort("created_at", -1).skip(skip).limit(limit)
    
    videos = await cursor.to_list(length=limit)
    return {"success": True, "videos": videos, "count": len(videos)}


@router.get("/concepts")
async def get_all_concepts(
    current_user: dict = Depends(get_current_user)
):
    """Get all RealApex demo concepts categorized"""
    check_admin(current_user)
    
    # This returns the concept structure for reference
    return {
        "success": True,
        "categories": [
            {
                "id": "proptech_foundation",
                "name": "PropTech Foundation",
                "concepts": ["Digital Property Management", "Smart Layout Mapping", "Multi-Project Portfolio", "Property Gallery & Tours", "Geo-Tagged Locations"]
            },
            {
                "id": "ai_sales_marketing", 
                "name": "AI-Powered Sales & Marketing",
                "concepts": ["AI Lead Scoring", "Automated Follow-ups", "Agent Performance Analytics", "Predictive Forecasting", "AI Video Generation", "Smart Segmentation"]
            },
            {
                "id": "customer_experience",
                "name": "Smart Customer Experience", 
                "concepts": ["Public Layout Portal", "Real-Time Status", "Self-Service Booking", "Payment Dashboard", "Site Visit Scheduling", "Document Center"]
            },
            {
                "id": "fintech_integration",
                "name": "FinTech & Payment Solutions",
                "concepts": ["Multi-Gateway Payments", "EMI Tracking", "Payment Reminders", "Bank Reconciliation", "Digital Receipts", "Revenue Analytics"]
            },
            {
                "id": "compliance_transparency",
                "name": "RERA Compliance & Transparency",
                "concepts": ["DLT SMS Templates", "Legal Document Management", "Booking Audit Trail", "KYC Verification", "Handover Documentation"]
            },
            {
                "id": "omnichannel_communication",
                "name": "Omnichannel Communication",
                "concepts": ["WhatsApp Integration", "SMS Campaigns", "Email Notifications", "Web Push", "In-App Notifications"]
            },
            {
                "id": "analytics_insights",
                "name": "Data Analytics & Insights",
                "concepts": ["Executive Dashboard", "Sales Pipeline", "Agent Reports", "Financial Health", "Customer Behavior", "Market Trends"]
            },
            {
                "id": "operational_efficiency",
                "name": "Operational Efficiency",
                "concepts": ["Role-Based Access", "Calendar Sync", "Workflow Automation", "Bulk Import/Export", "Multi-Branch Management"]
            }
        ]
    }


@router.post("/generate-voiceover")
async def generate_voiceover(
    request: GenerateVoiceoverRequest,
    current_user: dict = Depends(get_current_user)
):
    """Generate high-quality voiceover audio using Edge TTS (100% FREE - Unlimited)"""
    check_admin(current_user)
    
    try:
        import edge_tts
        import re
        
        # Edge TTS voice mapping
        EDGE_VOICES = {
            # English voices
            "nova": "en-US-JennyNeural",       # Female, friendly
            "alloy": "en-US-GuyNeural",        # Male, neutral
            "echo": "en-US-AriaNeural",        # Female, professional
            "fable": "en-US-DavisNeural",      # Male, narrator
            "onyx": "en-US-ChristopherNeural", # Male, deep
            "shimmer": "en-US-SaraNeural",     # Female, bright
            # Indian English
            "indian_male": "en-IN-PrabhatNeural",
            "indian_female": "en-IN-NeerjaNeural",
            # Telugu
            "telugu_male": "te-IN-MohanNeural",
            "telugu_female": "te-IN-ShrutiNeural",
            # Hindi
            "hindi_male": "hi-IN-MadhurNeural",
            "hindi_female": "hi-IN-SwaraNeural",
        }
        
        # Get voice or default
        voice = EDGE_VOICES.get(request.voice, "en-US-JennyNeural")
        
        # Clean script
        clean_script = re.sub(r'\[SCREEN:[^\]]*\]', '', request.script)
        clean_script = re.sub(r'\[PAUSE\]', '...', clean_script)
        clean_script = re.sub(r'\*\*([^*]+)\*\*', r'\1', clean_script)
        clean_script = re.sub(r'#{1,6}\s*', '', clean_script)
        clean_script = clean_script.strip()
        
        if not clean_script:
            raise HTTPException(status_code=400, detail="Script is empty")
        
        print("🎙️ Generating voiceover with Edge TTS (FREE)")
        print(f"📝 Script: {len(clean_script)} chars, Voice: {voice}")
        
        # Generate audio using Edge TTS
        filename = f"voiceover_{request.concept_title.replace(' ', '_').replace(':', '')}_{uuid.uuid4().hex[:8]}.mp3"
        filepath = f"/tmp/{filename}"
        
        # Calculate rate adjustment for Edge TTS
        if request.speed == 1.0:
            rate = "+0%"
        elif request.speed > 1:
            rate = f"+{int((request.speed - 1) * 100)}%"
        else:
            rate = f"-{int((1 - request.speed) * 100)}%"
        
        communicate = edge_tts.Communicate(clean_script, voice, rate=rate)
        await communicate.save(filepath)
        
        # Get file size
        file_size = os.path.getsize(filepath)
        
        print(f"✅ Voiceover generated: {filename} ({file_size} bytes)")
        
        return {
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "size_bytes": file_size,
            "voice_used": voice,
            "download_url": f"/api/realapex-demos/download-voiceover/{filename}"
        }
        
    except Exception as e:
        print(f"❌ Voiceover error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download-voiceover/{filename}")
async def download_voiceover(
    filename: str,
    current_user: dict = Depends(get_current_user)
):
    """Download generated voiceover MP3 file"""
    check_admin(current_user)
    
    filepath = f"/tmp/{filename}"
    
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Voiceover file not found")
    
    with open(filepath, "rb") as f:
        audio_bytes = f.read()
    
    return Response(
        content=audio_bytes,
        media_type="audio/mpeg",
        headers={
            "Content-Disposition": f"attachment; filename={filename}"
        }
    )


@router.get("/voice-options")
async def get_voice_options(
    current_user: dict = Depends(get_current_user)
):
    """Get available voice options for voiceover generation"""
    check_admin(current_user)
    
    return {
        "success": True,
        "voices": [
            {"id": "alloy", "name": "Alloy", "description": "Neutral, balanced - Good for formal content"},
            {"id": "ash", "name": "Ash", "description": "Clear, articulate - Good for tutorials"},
            {"id": "coral", "name": "Coral", "description": "Warm, friendly - Good for customer-facing"},
            {"id": "echo", "name": "Echo", "description": "Smooth, calm - Good for explainers"},
            {"id": "fable", "name": "Fable", "description": "Expressive, storytelling - Good for narratives"},
            {"id": "nova", "name": "Nova", "description": "Energetic, upbeat - Good for demos (Recommended)"},
            {"id": "onyx", "name": "Onyx", "description": "Deep, authoritative - Good for enterprise"},
            {"id": "sage", "name": "Sage", "description": "Wise, measured - Good for educational"},
            {"id": "shimmer", "name": "Shimmer", "description": "Bright, cheerful - Good for marketing"}
        ],
        "models": [
            {"id": "tts-1", "name": "Standard", "description": "Fast generation, good quality"},
            {"id": "tts-1-hd", "name": "HD Quality", "description": "Slower but higher quality (Recommended)"}
        ],
        "speed_range": {"min": 0.25, "max": 4.0, "default": 1.0}
    }



@router.post("/generate-presentation")
async def generate_presentation(
    request: GeneratePresentationRequest,
    current_user: dict = Depends(get_current_user)
):
    """Generate PowerPoint presentation from script (100% FREE - No AI/API needed)"""
    check_admin(current_user)
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        import re
        
        print(f"📊 Generating PPT (FREE) for: {request.concept_title}")
        
        # Step 1: Parse script into slides WITHOUT using any AI
        slides_data = parse_script_to_slides(request.script, request.concept_title, len(request.screenshot_urls))
        
        print(f"📝 Created {len(slides_data)} slides from script")
        
        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        # Theme colors based on selection
        themes = {
            "professional": {"bg": RGBColor(15, 23, 42), "accent": RGBColor(249, 115, 22), "text": RGBColor(255, 255, 255)},
            "modern": {"bg": RGBColor(30, 41, 59), "accent": RGBColor(168, 85, 247), "text": RGBColor(255, 255, 255)},
            "minimal": {"bg": RGBColor(255, 255, 255), "accent": RGBColor(59, 130, 246), "text": RGBColor(15, 23, 42)}
        }
        theme = themes.get(request.theme, themes["professional"])
        
        screenshot_index = 0
        
        for slide_info in slides_data:
            slide_type = slide_info.get("type", "content")
            
            # Add blank slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)
            
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme["bg"]
            
            if slide_type == "title":
                # Title slide
                title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_info.get("title", "")
                p.font.size = Pt(54)
                p.font.bold = True
                p.font.color.rgb = theme["text"]
                p.alignment = PP_ALIGN.CENTER
                
                if slide_info.get("subtitle"):
                    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(14), Inches(1))
                    tf2 = sub_box.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = slide_info.get("subtitle", "")
                    p2.font.size = Pt(28)
                    p2.font.color.rgb = theme["accent"]
                    p2.alignment = PP_ALIGN.CENTER
                    
            elif slide_type == "screenshot":
                # Screenshot placeholder slide
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_info.get("title", "Feature Demo")
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = theme["text"]
                
                # Placeholder for screenshot
                placeholder_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(5))
                tf2 = placeholder_box.text_frame
                p2 = tf2.paragraphs[0]
                p2.text = f"📸 Screenshot {screenshot_index + 1}\n\n{slide_info.get('caption', 'Add your app screenshot here')}"
                p2.font.size = Pt(24)
                p2.font.color.rgb = RGBColor(148, 163, 184)
                p2.alignment = PP_ALIGN.CENTER
                
                # Add actual screenshot if provided
                if request.screenshot_urls and screenshot_index < len(request.screenshot_urls):
                    # Note: For actual screenshots, would need to download and embed
                    pass
                
                screenshot_index += 1
                
            elif slide_type in ["content", "benefits"]:
                # Content slide with bullets
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_info.get("title", "")
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = theme["text"]
                
                # Bullets
                bullets = slide_info.get("bullets", [])
                if bullets:
                    bullet_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
                    tf2 = bullet_box.text_frame
                    tf2.word_wrap = True
                    
                    for i, bullet in enumerate(bullets):
                        if i == 0:
                            p2 = tf2.paragraphs[0]
                        else:
                            p2 = tf2.add_paragraph()
                        p2.text = f"• {bullet}"
                        p2.font.size = Pt(28)
                        p2.font.color.rgb = theme["text"]
                        p2.space_after = Pt(20)
                        
            elif slide_type == "cta":
                # Call to action slide
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(1.5))
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_info.get("title", "Get Started Today")
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = theme["accent"]
                p.alignment = PP_ALIGN.CENTER
                
                # CTA details
                bullets = slide_info.get("bullets", ["Visit RealApex.in", "Request a Free Demo"])
                cta_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(3))
                tf2 = cta_box.text_frame
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p2 = tf2.paragraphs[0]
                    else:
                        p2 = tf2.add_paragraph()
                    p2.text = bullet
                    p2.font.size = Pt(32)
                    p2.font.color.rgb = theme["text"]
                    p2.alignment = PP_ALIGN.CENTER
                    p2.space_after = Pt(15)
            
            # Add speaker notes
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_info.get("notes", "")
        
        # Save presentation
        filename = f"presentation_{request.concept_title.replace(' ', '_').replace(':', '')}_{uuid.uuid4().hex[:8]}.pptx"
        filepath = f"/tmp/{filename}"
        prs.save(filepath)
        
        file_size = os.path.getsize(filepath)
        
        print(f"✅ Presentation generated: {filename} ({file_size} bytes, {len(slides_data)} slides)")
        
        return {
            "success": True,
            "filename": filename,
            "filepath": filepath,
            "size_bytes": file_size,
            "slides_count": len(slides_data),
            "download_url": f"/api/realapex-demos/download-presentation/{filename}",
            "slides_structure": slides_data
        }
        
    except Exception as e:
        print(f"❌ Presentation generation error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download-presentation/{filename}")
async def download_presentation(
    filename: str,
    current_user: dict = Depends(get_current_user)
):
    """Download generated PowerPoint presentation"""
    check_admin(current_user)
    
    filepath = f"/tmp/{filename}"
    
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Presentation file not found")
    
    return FileResponse(
        filepath,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )


@router.post("/upload-screenshots")
async def upload_screenshots(
    files: List[UploadFile] = File(...),
    current_user: dict = Depends(get_current_user)
):
    """Upload screenshots for presentation"""
    check_admin(current_user)
    
    uploaded_files = []
    
    for file in files:
        if not file.content_type.startswith('image/'):
            continue
            
        # Save file
        file_id = uuid.uuid4().hex[:12]
        filename = f"screenshot_{file_id}_{file.filename}"
        filepath = f"/tmp/{filename}"
        
        content = await file.read()
        with open(filepath, "wb") as f:
            f.write(content)
        
        uploaded_files.append({
            "filename": filename,
            "filepath": filepath,
            "original_name": file.filename,
            "size": len(content),
            "url": f"/api/realapex-demos/screenshot/{filename}"
        })
    
    return {
        "success": True,
        "uploaded_count": len(uploaded_files),
        "files": uploaded_files
    }


@router.get("/screenshot/{filename}")
async def get_screenshot(filename: str):
    """Get uploaded/generated screenshot - public endpoint for image display"""
    # Made public so images can load in <img> tags without auth headers
    
    filepath = f"/tmp/{filename}"
    
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="Screenshot not found")
    
    # Determine content type
    if filename.lower().endswith('.png'):
        media_type = "image/png"
    elif filename.lower().endswith(('.jpg', '.jpeg')):
        media_type = "image/jpeg"
    else:
        media_type = "image/png"
    
    return FileResponse(filepath, media_type=media_type)



class AutoGenerateImagesRequest(BaseModel):
    concept_title: str
    script: Optional[str] = ""
    category_name: Optional[str] = ""
    num_images: int = 3


@router.post("/auto-generate-images")
async def auto_generate_images(
    request: AutoGenerateImagesRequest,
    current_user: dict = Depends(get_current_user)
):
    """Auto-generate relevant images using OpenAI Image Generation (FREE with Emergent Key)"""
    check_admin(current_user)
    
    try:
        from emergentintegrations.llm.openai.image_generation import OpenAIImageGeneration
        
        EMERGENT_LLM_KEY = os.getenv('EMERGENT_LLM_KEY', '')
        
        if not EMERGENT_LLM_KEY:
            raise HTTPException(status_code=500, detail="EMERGENT_LLM_KEY not configured")
        
        # Generate image prompts based on concept
        image_prompts = []
        
        # Create relevant prompts based on the concept
        base_context = "Professional SaaS software dashboard screenshot, clean modern UI design, dark theme"
        
        if "layout" in request.concept_title.lower() or "plot" in request.concept_title.lower():
            image_prompts = [
                f"{base_context}, real estate property layout map with colored plots showing available (green), booked (yellow), and sold (red) status, interactive grid view",
                f"{base_context}, property management dashboard showing plot details, pricing information, customer data in a sidebar panel",
                f"{base_context}, mobile-responsive real estate app showing property gallery with images and status indicators"
            ]
        elif "payment" in request.concept_title.lower() or "emi" in request.concept_title.lower() or "fintech" in request.concept_title.lower():
            image_prompts = [
                f"{base_context}, payment dashboard showing EMI schedule, due dates, payment history with charts and graphs",
                f"{base_context}, financial analytics dashboard with revenue graphs, collection reports, payment gateway integration",
                f"{base_context}, customer payment portal showing balance, upcoming payments, receipt download options"
            ]
        elif "lead" in request.concept_title.lower() or "sales" in request.concept_title.lower() or "crm" in request.concept_title.lower():
            image_prompts = [
                f"{base_context}, CRM lead management dashboard showing lead pipeline, conversion funnel, priority scores",
                f"{base_context}, sales analytics dashboard with charts showing lead sources, conversion rates, agent performance",
                f"{base_context}, customer profile page showing contact details, interaction history, follow-up reminders"
            ]
        elif "analytics" in request.concept_title.lower() or "dashboard" in request.concept_title.lower() or "report" in request.concept_title.lower():
            image_prompts = [
                f"{base_context}, executive dashboard with KPI cards, revenue charts, sales metrics, performance indicators",
                f"{base_context}, analytics page with pie charts, bar graphs, trend lines showing business performance",
                f"{base_context}, report generation interface with filters, date range selectors, export options"
            ]
        elif "sms" in request.concept_title.lower() or "whatsapp" in request.concept_title.lower() or "notification" in request.concept_title.lower():
            image_prompts = [
                f"{base_context}, notification center showing SMS templates, WhatsApp messages, delivery status",
                f"{base_context}, messaging dashboard with template editor, recipient list, send history",
                f"{base_context}, communication settings page with channel configuration, automation rules"
            ]
        else:
            # Generic SaaS dashboard prompts
            image_prompts = [
                f"{base_context}, main dashboard overview with key metrics cards, quick action buttons, recent activity feed",
                f"{base_context}, feature showcase screen with step-by-step workflow, highlighted UI elements",
                f"{base_context}, settings and configuration page with form inputs, toggle switches, save buttons"
            ]
        
        # Limit to requested number
        image_prompts = image_prompts[:request.num_images]
        
        # Generate images
        image_gen = OpenAIImageGeneration(api_key=EMERGENT_LLM_KEY)
        
        generated_images = []
        
        for i, prompt in enumerate(image_prompts):
            print(f"🎨 Generating image {i+1}/{len(image_prompts)}")
            
            try:
                # generate_images returns List[bytes]
                image_bytes_list = await image_gen.generate_images(
                    prompt=prompt,
                    model="gpt-image-1",
                    number_of_images=1,
                    quality="low"  # Use low for faster generation
                )
                
                if image_bytes_list and len(image_bytes_list) > 0:
                    # Save image to file
                    img_filename = f"generated_img_{uuid.uuid4().hex[:8]}.png"
                    img_filepath = f"/tmp/{img_filename}"
                    
                    with open(img_filepath, "wb") as f:
                        f.write(image_bytes_list[0])
                    
                    generated_images.append({
                        "url": f"/api/realapex-demos/screenshot/{img_filename}",
                        "filename": img_filename,
                        "prompt": prompt[:100],
                        "index": i + 1
                    })
                    print(f"✅ Image {i+1} generated and saved: {img_filename}")
            except Exception as img_error:
                print(f"⚠️ Image {i+1} failed: {str(img_error)}")
                import traceback
                traceback.print_exc()
                continue
        
        return {
            "success": True,
            "images": generated_images,
            "total_requested": request.num_images,
            "total_generated": len(generated_images)
        }
        
    except Exception as e:
        print(f"❌ Auto-generate images error: {str(e)}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ================== YOUTUBE CONTENT GENERATOR (AgentApex) ==================

YOUTUBE_CONTENT_CATEGORIES = [
    {
        "id": "property_tips",
        "name": "Property Tips",
        "icon": "💡",
        "description": "Buying, selling, and investment tips"
    },
    {
        "id": "area_reviews",
        "name": "Area Reviews",
        "icon": "📍",
        "description": "Location analysis and reviews"
    },
    {
        "id": "market_updates",
        "name": "Market Updates",
        "icon": "📈",
        "description": "Current market trends and prices"
    },
    {
        "id": "investment_guide",
        "name": "Investment Guide",
        "icon": "💰",
        "description": "Investment strategies and returns"
    },
    {
        "id": "legal_tips",
        "name": "Legal Tips",
        "icon": "⚖️",
        "description": "Documentation and legal advice"
    },
    {
        "id": "success_stories",
        "name": "Success Stories",
        "icon": "🏆",
        "description": "Client testimonials and case studies"
    }
]

class YouTubeContentRequest(BaseModel):
    topic: str
    category: str
    language: str = "english"
    tone: str = "professional"  # professional, friendly, motivational
    target_audience: str = "property_buyers"
    include_emotional: bool = True
    custom_context: Optional[str] = ""

class YouTubeContentHistory(BaseModel):
    id: str
    topic: str
    category: str
    content: str
    language: str
    created_at: str
    published: bool = False
    seo_slug: Optional[str] = None

@router.get("/youtube-content/categories")
async def get_youtube_categories():
    """Get available YouTube content categories"""
    return {"categories": YOUTUBE_CONTENT_CATEGORIES}

@router.post("/youtube-content/generate")
async def generate_youtube_content(request: YouTubeContentRequest):
    """Generate YouTube script content using Claude AI - with emotional intelligence"""
    try:
        from emergentintegrations.llm.chat import LlmChat, UserMessage
    except ImportError:
        raise HTTPException(status_code=503, detail="AI service unavailable")
    
    api_key = os.environ.get('EMERGENT_LLM_KEY')
    if not api_key:
        raise HTTPException(status_code=503, detail="AI service not configured")
    
    # Build emotional tone guidance
    tone_guide = {
        "professional": "Use confident, authoritative language. Build trust through expertise.",
        "friendly": "Use warm, approachable language. Connect personally with viewers like a helpful friend.",
        "motivational": "Use inspiring, energetic language. Encourage action and build excitement."
    }
    
    emotional_guidance = ""
    if request.include_emotional:
        emotional_guidance = """
EMOTIONAL INTELLIGENCE GUIDELINES:
- Start with a relatable problem or aspiration the viewer has
- Use storytelling elements - paint scenarios they can visualize
- Include moments of empathy ("I know how frustrating it can be when...")
- Build confidence gradually ("Here's the good news...")
- End with hope and clear next steps
- Use power words: discover, transform, secure, achieve, protect, thrive
"""
    
    category_info = next((c for c in YOUTUBE_CONTENT_CATEGORIES if c["id"] == request.category), None)
    category_name = category_info["name"] if category_info else request.category
    
    prompt = f"""You are an expert real estate content creator for YouTube in India. Create COPY-PASTE READY content for a video.

TOPIC: {request.topic}
CATEGORY: {category_name}
LANGUAGE: {request.language}
TONE: {tone_guide.get(request.tone, tone_guide['professional'])}
TARGET AUDIENCE: {request.target_audience}

{emotional_guidance}

{f"ADDITIONAL CONTEXT: {request.custom_context}" if request.custom_context else ""}

CREATE THE FOLLOWING (ready to copy-paste):

1. **VIDEO TITLE** (catchy, SEO-friendly, max 60 chars)

2. **VIDEO DESCRIPTION** (for YouTube description box, include relevant keywords)

3. **HASHTAGS** (10-15 relevant hashtags)

4. **FULL SCRIPT/CONTENT**
Write the complete spoken content as if speaking directly to the viewer.
- DO NOT include timestamps or duration markers
- DO NOT include [INTRO], [OUTRO] or any section markers
- Just write the natural, flowing script that can be read directly
- Include hook, main points, and strong call-to-action
- Make it conversational and engaging

5. **SEO KEYWORDS** (comma separated, for internal use)

6. **THUMBNAIL TEXT SUGGESTION** (2-4 words max for overlay text)

IMPORTANT:
- Content should be 100% ready to copy-paste
- No instructions or meta-comments
- No time markers (like "0:00-0:30")
- Just pure, usable content
"""
    
    try:
        chat = LlmChat(
            api_key=api_key,
            session_id=f"yt-content-{uuid.uuid4()}",
            system_message="You are a professional YouTube content creator specializing in Indian real estate. Create engaging, emotional, and actionable content."
        ).with_model("anthropic", "claude-sonnet-4-20250514")
        
        user_message = UserMessage(text=prompt)
        response = await chat.send_message(user_message)
        
        # Store in database
        from motor.motor_asyncio import AsyncIOMotorClient
        from datetime import datetime, timezone
        
        content_id = str(uuid.uuid4())
        mongo_url = os.environ.get('MONGO_URL')
        if mongo_url:
            client = AsyncIOMotorClient(mongo_url)
            db = client[os.environ.get('DB_NAME', 'realapex')]
            
            content_record = {
                "id": content_id,
                "topic": request.topic,
                "category": request.category,
                "language": request.language,
                "tone": request.tone,
                "content": response,
                "published": False,
                "seo_slug": None,
                "created_at": datetime.now(timezone.utc).isoformat()
            }
            
            await db.youtube_content_history.insert_one(content_record)
            client.close()
        
        return {
            "success": True,
            "content": response,
            "content_id": content_id,
            "topic": request.topic,
            "category": category_name
        }
        
    except Exception as e:
        print(f"YouTube content generation error: {e}")
        raise HTTPException(status_code=500, detail=f"Content generation failed: {str(e)}")

@router.get("/youtube-content/history")
async def get_youtube_content_history():
    """Get history of generated YouTube content"""
    from motor.motor_asyncio import AsyncIOMotorClient
    
    mongo_url = os.environ.get('MONGO_URL')
    if not mongo_url:
        return {"history": []}
    
    try:
        client = AsyncIOMotorClient(mongo_url)
        db = client[os.environ.get('DB_NAME', 'realapex')]
        
        history = await db.youtube_content_history.find(
            {}, {"_id": 0}
        ).sort("created_at", -1).limit(50).to_list(50)
        
        client.close()
        return {"history": history}
    except Exception as e:
        print(f"Error fetching history: {e}")
        return {"history": []}

@router.put("/youtube-content/{content_id}/publish")
async def publish_content_as_seo(content_id: str, seo_slug: str):
    """Mark content as published and set SEO slug for website"""
    from motor.motor_asyncio import AsyncIOMotorClient
    
    mongo_url = os.environ.get('MONGO_URL')
    if not mongo_url:
        raise HTTPException(status_code=503, detail="Database not configured")
    
    try:
        client = AsyncIOMotorClient(mongo_url)
        db = client[os.environ.get('DB_NAME', 'realapex')]
        
        result = await db.youtube_content_history.update_one(
            {"id": content_id},
            {"$set": {"published": True, "seo_slug": seo_slug}}
        )
        
        client.close()
        
        if result.modified_count == 0:
            raise HTTPException(status_code=404, detail="Content not found")
        
        return {"success": True, "message": "Content marked as published", "seo_slug": seo_slug}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.delete("/youtube-content/{content_id}")
async def delete_youtube_content(content_id: str):
    """Delete a YouTube content record"""
    from motor.motor_asyncio import AsyncIOMotorClient
    
    mongo_url = os.environ.get('MONGO_URL')
    if not mongo_url:
        raise HTTPException(status_code=503, detail="Database not configured")
    
    try:
        client = AsyncIOMotorClient(mongo_url)
        db = client[os.environ.get('DB_NAME', 'realapex')]
        
        result = await db.youtube_content_history.delete_one({"id": content_id})
        client.close()
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Content not found")
        
        return {"success": True, "message": "Content deleted"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))



# ===================== PUBLIC SEO PAGES =====================

@router.get("/seo-articles")
async def get_published_articles():
    """Get all published articles for SEO listing page"""
    from motor.motor_asyncio import AsyncIOMotorClient
    
    mongo_url = os.environ.get('MONGO_URL')
    if not mongo_url:
        return {"articles": []}
    
    try:
        client = AsyncIOMotorClient(mongo_url)
        db = client[os.environ.get('DB_NAME', 'realapex')]
        
        articles = await db.youtube_content_history.find(
            {"published": True},
            {"_id": 0, "id": 1, "topic": 1, "category": 1, "language": 1, "seo_slug": 1, "created_at": 1}
        ).sort("created_at", -1).to_list(100)
        
        client.close()
        return {"articles": articles}
    except Exception as e:
        print(f"Error fetching articles: {e}")
        return {"articles": []}

@router.get("/seo-article/{slug}")
async def get_seo_article(slug: str):
    """Get a single published article by SEO slug for public viewing"""
    from motor.motor_asyncio import AsyncIOMotorClient
    
    mongo_url = os.environ.get('MONGO_URL')
    if not mongo_url:
        raise HTTPException(status_code=503, detail="Database not configured")
    
    try:
        client = AsyncIOMotorClient(mongo_url)
        db = client[os.environ.get('DB_NAME', 'realapex')]
        
        article = await db.youtube_content_history.find_one(
            {"seo_slug": slug, "published": True},
            {"_id": 0}
        )
        
        client.close()
        
        if not article:
            raise HTTPException(status_code=404, detail="Article not found")
        
        return {"success": True, "article": article}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
